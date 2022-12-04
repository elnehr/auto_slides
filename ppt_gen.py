import openai
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt

# Set up the GPT-3 API
openai.api_key_path = "OPENAI_API_KEY.txt"


topic = input("Enter the topic: ")

# Generate text for the slide using the GPT-3 API
prompt = (
    "Please generate text for a McKinsey-style slide about the pros and cons of " + topic +". "
    "The slide should include a title, a list of pros, a list of cons, and a key takeaway. "
    "The text should be formatted as follows: "
    "Title: [insert title here, remember to use top down communication] "
    "Pros: "
    "Pros: - [insert pro 1] "
    "Pros: - [insert pro 2] "
    "Pros: - [insert pro 3] "
    "Pros: - [insert pro 4] "
    "Pros: - [insert pro 5] "
    "Cons: "
    "Cons: - [insert con 1] "
    "Cons: - [insert con 2] "
    "Cons: - [insert con 3] "
    "Cons: - [insert con 4] "
    "Cons: - [insert con 5] "
    "Key takeaway: "
    "- [insert key takeaway]"
)
completions = openai.Completion.create(
    engine="text-davinci-003",
    prompt=prompt,
    temperature=0.7,
    max_tokens=300,
    top_p=1,
    frequency_penalty=0,
    presence_penalty=0
)


generated_text = completions.choices[0].text
lines = generated_text.split("\n")
title = lines[2][7:]
pros = lines[4:9]
cons = lines[10:15]
key_takeaway = None
current_list = None
for line in lines[1:]:
    if line.strip().startswith("Key Takeaway:"):
        key_takeaway = line.strip()[14:]
        continue
        
for line in lines[1:]:
    line = line.strip()[1:]
        
        
pros = lines[4:9]
cons = lines[10:15]


prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = title

tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = 'Pros'
p.font.size = Pt(16)

for pro in pros:
    pro= pro.strip()[1:]
    p = tf.add_paragraph()
    p.text = pro
    p.font.size = Pt(16)
    p.level = 1
 

p = tf.add_paragraph()
p.text = 'Cons'
p.font.size = Pt(16)

for con in cons:
    con = con.strip()[1:]
    p = tf.add_paragraph()
    p.text = con
    p.font.size = Pt(16)
    p.level = 1

p = tf.add_paragraph()
p.text = key_takeaway
p.font.size = Pt(16)


prs.save('test.pptx')